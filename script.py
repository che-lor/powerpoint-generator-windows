#!/usr/bin/python3

from pptx import Presentation
from Fill_Out import *
import itertools

files = [powerpoint_name, file1, file2, file3]
ext = [".pptx", ".txt"]
file_dir = r"C:\Users\Home\Desktop\powerpoint-generator\Songs\\"

#Store fullname of files to list
file_names = []
for f in files:
	if f == powerpoint_name:
		file_names.append(f + ext[0])
	else:
		file_names.append(f + ext[1])

#Create new list for slides
lines = []
title_cycle = itertools.cycle(files[1:])
for fn in file_names[1:]:
	file = open((file_dir + fn), 'r')
	file_lines = file.read().splitlines()
	file.close()
	#Cycle through file_names list and append them to lines list
	lines.append(next(title_cycle))
	#Append blank line after
	lines.append("")
	#Extend the file_lines after
	lines.extend(file_lines)
	#Add extra line after if lines are odd
	if ((int(len(file_lines))) % LPS) != 0:
		lines.append("")
		print("Added extra line to " + fn)
	#Add lines to create blank slide between each file
	for n in range(LPS):
		lines.append("")

#Open presentation using the "slide_master.pptx" file
prs = Presentation("slide_master.pptx")

#Create slides for each LPS
total_slides = int(len(lines)/LPS)
line_cycle = itertools.cycle(lines)
loop = 0
while loop < total_slides:
	pair = next(line_cycle) + '\n' + next(line_cycle) # add another cycle for more LPS
	bullet_slide_layout = prs.slide_layouts[1]
	slide = prs.slides.add_slide(bullet_slide_layout)
	shapes = slide.shapes
	body_shape = shapes.placeholders[1]
	tf = body_shape.text_frame
	tf.text = pair
	loop += 1
	if loop > total_slides:
		break

#Save file
save_dir = r"C:\Users\Home\Desktop\powerpoint-generator\Presentations\\"
save = prs.save(save_dir + file_names[0])
if [save]:
	print(file_names[0] + " was created!")